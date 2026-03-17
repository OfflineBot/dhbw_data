#!/usr/bin/env python3
"""Erstellt airline_satisfaction_praesi.pptx aus der Vorlage."""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
import copy

DARK = RGBColor(0x20, 0x20, 0x20)
GRAY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation("origins/Vorlage.pptx")

CONTENT_TOP  = 1_635_646
SUBTITLE_TOP = 1_131_590


# ── Helpers ───────────────────────────────────────────────────────────────────

def shape_near(slide, top_target, tol=300_000):
    best, best_d = None, float("inf")
    for s in slide.shapes:
        d = abs(s.top - top_target)
        if d < tol and d < best_d:
            best_d, best = d, s
    return best


def set_text(shape, lines, size=13, color=DARK):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_img(slide, path, left, top, w, h):
    slide.shapes.add_picture(path, Emu(left), Emu(top), Emu(w), Emu(h))


def fix_date(slide):
    for s in slide.shapes:
        if s.has_text_frame and "Q1 - 2025" in s.text_frame.text:
            set_text(s, ["Q1 - 2026"], size=10, color=DARK)


def set_pagenum(slide, num):
    for s in slide.shapes:
        if s.has_text_frame and s.left > 8_000_000 and s.top > 4_500_000:
            set_text(s, [str(num)], size=10, color=DARK)


def dup_slide(prs, idx):
    """Dupliziert Folie an idx, hängt sie ans Ende."""
    src = prs.slides[idx]
    new = prs.slides.add_slide(src.slide_layout)
    src_tree = src.shapes._spTree
    new_tree = new.shapes._spTree
    for el in list(new_tree):
        new_tree.remove(el)
    for el in src_tree:
        new_tree.append(copy.deepcopy(el))
    return new


def reorder(prs, order):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for el in list(lst):
        lst.remove(el)
    for i in order:
        lst.append(ids[i])


# ── Exploration-Folie 4× duplizieren (landet als idx 7–10) ────────────────────
for _ in range(4):
    dup_slide(prs, 4)

# Gewünschte Reihenfolge: Titel, Agenda, Anwendungsfall, Cleaning,
#                         Exploration ×5, Modelling, Ausblick
reorder(prs, [0, 1, 2, 3, 4, 7, 8, 9, 10, 5, 6])

# Referenz auf alle Folien (nach Reorder)
sl = prs.slides   # 0..10

# ── Slide 0 · Titelfolie ──────────────────────────────────────────────────────
for shape in sl[0].shapes:
    if not shape.has_text_frame:
        continue
    txt = shape.text_frame.text
    if "Vorname" in txt or "Data Visualization" in txt:
        set_text(shape, ["Airline Passenger Satisfaction", "von Leon"], size=24, color=DARK)
        shape.text_frame.paragraphs[0].runs[0].font.bold = True
    elif "03.2025" in txt:
        set_text(shape, ["03.2026"], size=12, color=DARK)
    elif "Posterpräsentation" in txt:
        set_text(shape, ["WDS125 – Data Visualization"], size=14, color=GRAY)

# ── Slide 1 · Gliederung ──────────────────────────────────────────────────────
fix_date(sl[1])
set_pagenum(sl[1], 2)

# ── Slide 2 · Anwendungsfall ──────────────────────────────────────────────────
fix_date(sl[2]); set_pagenum(sl[2], 3)

shape_near(sl[2], SUBTITLE_TOP) and set_text(
    shape_near(sl[2], SUBTITLE_TOP),
    ["103.902 Passagiere · 11 Features · Zielvariable: Zufriedenheit"],
    size=12, color=GRAY
)
c = shape_near(sl[2], CONTENT_TOP)
if c:
    c.width = Emu(3_900_000)
    set_text(c, [
        "satisfied         → 43,3 %",
        "not satisfied  → 56,7 %",
    ])
add_img(sl[2], "figures/fig_target_distribution.png",
        4_550_000, 1_450_000, 4_250_000, 3_200_000)

# ── Slide 3 · Data Cleaning ───────────────────────────────────────────────────
fix_date(sl[3]); set_pagenum(sl[3], 4)

shape_near(sl[3], SUBTITLE_TOP) and set_text(
    shape_near(sl[3], SUBTITLE_TOP),
    ["Bereinigte Datenpunkte vor der Analyse"],
    size=12, color=GRAY
)
c = shape_near(sl[3], CONTENT_TOP)
if c:
    c.width = Emu(3_900_000)
    set_text(c, [
        "2 ungültige Leg room service Werte entfernt",
        "3 fehlende Age-Werte entfernt",
        "Spalten Customer ID & isPilot gedroppt",
        "",
        "Ausreißer (Age –24, Distance –200)",
        "→ sichtbar, aber nicht bereinigt",
    ])
add_img(sl[3], "figures/fig_outlier_boxplot.png",
        4_550_000, 1_450_000, 4_250_000, 3_200_000)

# ── Exploration-Folien 4–8 ────────────────────────────────────────────────────
exp_slides = [
    (4,  "Flugklassen",
     ["Business klar dominierend",
      "Eco mit dem kleinsten Anteil"],
     "figures/fig_class_distribution.png"),

    (5,  "Altersverteilung",
     ["Schwerpunkt: 25–50 Jahre",
      "Kaum Passagiere unter 15 oder über 80"],
     "figures/fig_age_histogram.png"),

    (6,  "Flugdistanz nach Zufriedenheit",
     ["Zufriedene Passagiere fliegen im Median weiter",
      "Große Streuung in beiden Gruppen"],
     "figures/fig_distance_boxplot.png"),

    (7,  "Korrelationsmatrix",
     ["Food & drink ↔ On-board service: r ≈ 0.44",
      "Departure Delay kaum mit anderen Features korreliert"],
     "figures/fig_correlation_heatmap.png"),

    (8,  "Alter vs. Flugdistanz",
     ["Kein klarer Zusammenhang erkennbar",
      "Zufriedenheit nicht durch diese beiden Features erklärbar"],
     "figures/fig_scatter.png"),
]

for idx, subtitle, bullets, img in exp_slides:
    s = sl[idx]
    fix_date(s); set_pagenum(s, idx + 1)
    sub = shape_near(s, SUBTITLE_TOP)
    if sub:
        set_text(sub, [subtitle], size=12, color=GRAY)
    c = shape_near(s, CONTENT_TOP)
    if c:
        c.width = Emu(3_900_000)
        set_text(c, bullets)
    add_img(s, img, 4_550_000, 1_450_000, 4_250_000, 3_200_000)

# ── Slide 9 · Predictive Modelling ───────────────────────────────────────────
fix_date(sl[9]); set_pagenum(sl[9], 10)

shape_near(sl[9], SUBTITLE_TOP) and set_text(
    shape_near(sl[9], SUBTITLE_TOP),
    ["RandomForestClassifier · Split 80/10/10"],
    size=12, color=GRAY
)
c = shape_near(sl[9], CONTENT_TOP)
if c:
    c.width = Emu(3_900_000)
    set_text(c, [
        "Accuracy:  86 %",
        "",
        "dissatisfied  →  Recall 92 %",
        "satisfied       →  Recall 79 %",
        "",
        "Unzufriedene werden besser erkannt",
    ])
add_img(sl[9], "figures/fig_confusion_matrix.png",
        4_700_000, 1_450_000, 3_850_000, 3_200_000)

# ── Slide 10 · Nächste Schritte ───────────────────────────────────────────────
fix_date(sl[10]); set_pagenum(sl[10], 11)

shape_near(sl[10], SUBTITLE_TOP) and set_text(
    shape_near(sl[10], SUBTITLE_TOP),
    ["Verbesserungspotenzial bei den Visualisierungen"],
    size=12, color=GRAY
)
c = shape_near(sl[10], CONTENT_TOP)
if c:
    set_text(c, [
        "Scatter: 100k+ Punkte → Sampling oder Hexbin-Plot",
        "",
        "Konfusionsmatrix: Prozentanzeige statt Absolutwerte",
        "",
        "Interaktive Plots mit Plotly:",
        "  Filter nach Klasse, Hover-Infos, Zoom",
        "",
        "Farbpalette: colorblind-safe (z.B. viridis)",
        "",
        "Histogramm: KDE-Kurve ergänzen",
    ])

# ── Speichern ─────────────────────────────────────────────────────────────────
prs.save("airline_satisfaction_praesi.pptx")
print("Gespeichert: airline_satisfaction_praesi.pptx  (11 Folien)")
